"""
Convert Chinese course PowerPoint presentations into AI-friendly Markdown files.
Each PPTX becomes one MD file per lesson, structured for drill generation.
"""

import os
import re
import subprocess
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
INPUT_DIR = os.path.join(BASE_DIR, "resources", "in_person_lessons")
OUTPUT_DIR = os.path.join(BASE_DIR, "lessons_md")


def get_version() -> str:
    """Get version string from git hash + current date."""
    try:
        git_hash = subprocess.check_output(
            ["git", "rev-parse", "--short", "HEAD"],
            cwd=BASE_DIR, stderr=subprocess.DEVNULL
        ).decode().strip()
    except Exception:
        git_hash = "nogit"
    date = datetime.now().strftime("%Y-%m-%d %H:%M")
    return f"v{git_hash} ({date})"


def extract_slide(slide, slide_num: int) -> dict:
    """Extract all text content and tables from a single slide."""
    texts = []
    tables = []
    notes = ""

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)

        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_cells = [cell.text.strip() for cell in row.cells]
                table_data.append(row_cells)
            tables.append(table_data)

    # Extract speaker notes if present
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    return {
        "slide_num": slide_num,
        "texts": texts,
        "tables": tables,
        "notes": notes,
    }


def detect_lesson_number(filename: str, first_slide_texts: list[str]) -> str:
    """Try to determine lesson number from filename or first slide."""
    # Check filename
    m = re.search(r'(\d+)', filename)
    # Check first slide text for Roman or Arabic numerals
    for t in first_slide_texts:
        rm = re.search(r'Chinese\s+([IVX]+\.?)', t, re.IGNORECASE)
        if rm:
            roman = rm.group(1).rstrip('.')
            roman_map = {'I': 1, 'II': 2, 'III': 3, 'IV': 4, 'V': 5,
                         'VI': 6, 'VII': 7, 'VIII': 8, 'IX': 9, 'X': 10}
            if roman in roman_map:
                return str(roman_map[roman])
        am = re.search(r'Chinese\s+(\d+)', t, re.IGNORECASE)
        if am:
            return am.group(1)

    if m:
        return m.group(1)
    return "unknown"


def detect_lesson_title(first_slide_texts: list[str]) -> str:
    """Extract lesson title from first slide."""
    for t in first_slide_texts:
        if 'chinese' in t.lower() or 'hodina' in t.lower():
            return t
    return first_slide_texts[0] if first_slide_texts else "Untitled Lesson"


def is_pronunciation_pair(text: str) -> bool:
    """Check if a line looks like a pronunciation mapping (e.g. 'B=P' or 'Bian = Pien')."""
    return bool(re.match(r'^[A-Za-zÀ-ž]+\s*=\s*[A-Za-zÀ-ž()\s]+$', text))


def format_table_as_md(table_data: list[list[str]]) -> str:
    """Convert table data to markdown table."""
    if not table_data:
        return ""
    lines = []
    # Header row
    lines.append("| " + " | ".join(table_data[0]) + " |")
    lines.append("| " + " | ".join(["---"] * len(table_data[0])) + " |")
    for row in table_data[1:]:
        # Pad row if needed
        while len(row) < len(table_data[0]):
            row.append("")
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def categorize_content(slides_data: list[dict]) -> dict:
    """Organize slide content into semantic categories for AI-friendly output."""
    categories = {
        "overview": [],
        "phonetics": [],       # initials, finals, tones, pinyin
        "vocabulary": [],      # words, meanings
        "grammar": [],         # grammar rules, sentence patterns
        "writing": [],         # characters, stroke order
        "exercises": [],       # practice items, homework
        "resources": [],       # links, references
        "other": [],
    }

    phonetics_keywords = [
        'iniciál', 'finál', 'tón', 'pinyin', 'slabik', 'výslovnosť',
        'spoluhlásk', 'samohlásk', 'initial', 'final', 'tone',
        'pronunciation', 'fonetik',
    ]
    vocab_keywords = ['slovn', 'vocab', 'slov', 'hsk', 'učebnic', 'čítani']
    grammar_keywords = ['gramatik', 'grammar', 'slovosled', 'veta', 'sentence', 'izolačn']
    writing_keywords = ['písmo', 'znak', 'character', 'writing', 'stroke', 'logograf']
    exercise_keywords = ['domác', 'príprav', 'homework', 'cvičen', 'exercise', 'aktivit']
    resource_keywords = ['http', 'www', 'odkaz', 'link', 'slovník', 'pleco']

    for sd in slides_data:
        all_text = " ".join(sd["texts"]).lower()

        if sd["slide_num"] == 1:
            categories["overview"].append(sd)
        elif any(k in all_text for k in exercise_keywords):
            categories["exercises"].append(sd)
        elif any(k in all_text for k in resource_keywords):
            categories["resources"].append(sd)
        elif any(k in all_text for k in phonetics_keywords) or _has_pronunciation_pairs(sd["texts"]):
            categories["phonetics"].append(sd)
        elif any(k in all_text for k in vocab_keywords):
            categories["vocabulary"].append(sd)
        elif any(k in all_text for k in grammar_keywords):
            categories["grammar"].append(sd)
        elif any(k in all_text for k in writing_keywords):
            categories["writing"].append(sd)
        else:
            # Try to auto-detect based on content patterns
            if _has_chinese_chars(all_text) and _has_pinyin(all_text):
                categories["vocabulary"].append(sd)
            elif _has_pronunciation_pairs(sd["texts"]):
                categories["phonetics"].append(sd)
            else:
                categories["other"].append(sd)

    return categories


def _has_chinese_chars(text: str) -> bool:
    return bool(re.search(r'[\u4e00-\u9fff]', text))


def _has_pinyin(text: str) -> bool:
    pinyin_vowels = 'āáǎàēéěèīíǐìōóǒòūúǔùǖǘǚǜ'
    return any(c in text for c in pinyin_vowels)


def _has_pronunciation_pairs(texts: list[str]) -> bool:
    pairs = sum(1 for t in texts if is_pronunciation_pair(t))
    return pairs >= 2


def _join_fragmented_pairs(texts: list[str]) -> list[str]:
    """Join fragmented pronunciation items like ['P', 'Bian', '=', 'Pien'] into 'Bian = Pien (P)'.
    Also handles tab-separated pairs like 'Pian = Pchien\\t\\t\\tLian= Lien'."""
    result = []

    # First, split any tab-separated compound lines
    expanded = []
    for t in texts:
        if '\t' in t:
            parts = re.split(r'\t+', t)
            expanded.extend(p.strip() for p in parts if p.strip())
        else:
            expanded.append(t)

    # Detect the pattern: SK_consonant, Pinyin, =, Slovak_reading (groups of 4)
    i = 0
    while i < len(expanded):
        # Check for pattern: [consonant_label] [pinyin] [=] [reading]
        if (i + 3 < len(expanded)
                and expanded[i + 2] == '='
                and len(expanded[i]) <= 6
                and not re.match(r'^https?://', expanded[i])):
            label = expanded[i]       # e.g. "P" or "Pch"
            pinyin = expanded[i + 1]  # e.g. "Bian"
            reading = expanded[i + 3] # e.g. "Pien"
            result.append(f"{pinyin} = {reading} ({label})")
            i += 4
        else:
            # Try splitting "X = Y" that are already joined
            result.append(expanded[i])
            i += 1

    return result


def render_slide_content(sd: dict, heading_level: int = 3) -> str:
    """Render a single slide's content as markdown."""
    lines = []

    # Pre-process: join fragmented pronunciation pairs
    texts = _join_fragmented_pairs(sd["texts"])
    if not texts:
        return ""

    start_idx = 0
    if len(texts) > 1 and len(texts[0]) < 60 and not is_pronunciation_pair(texts[0]):
        lines.append(f"{'#' * heading_level} {texts[0]}")
        lines.append("")
        start_idx = 1

    # Check if content is mostly pronunciation pairs -> render as table
    remaining = texts[start_idx:]
    pron_pairs = [(t, is_pronunciation_pair(t)) for t in remaining]
    pair_count = sum(1 for _, is_p in pron_pairs if is_p)

    if pair_count > len(remaining) * 0.4 and pair_count >= 2:
        # Render as pronunciation table
        lines.append("| Pinyin | Slovak Pronunciation |")
        lines.append("| --- | --- |")
        for t, is_p in pron_pairs:
            if is_p:
                parts = re.split(r'\s*=\s*', t, maxsplit=1)
                if len(parts) == 2:
                    lines.append(f"| {parts[0].strip()} | {parts[1].strip()} |")
            else:
                # Non-pair line, add as note
                if t.strip():
                    lines.append("")
                    lines.append(f"> {t}")
        lines.append("")
    else:
        for t in remaining:
            # Detect Chinese characters with pinyin/translation
            if _has_chinese_chars(t):
                lines.append(f"- **{t}**")
            elif re.match(r'^https?://', t):
                lines.append(f"- <{t}>")
            else:
                lines.append(f"- {t}")
        lines.append("")

    # Render tables
    for table in sd["tables"]:
        lines.append(format_table_as_md(table))
        lines.append("")

    # Render notes
    if sd["notes"]:
        lines.append(f"> **Note:** {sd['notes']}")
        lines.append("")

    return "\n".join(lines)


CATEGORY_TITLES = {
    "overview": "Lesson Overview",
    "phonetics": "Phonetics & Pronunciation",
    "vocabulary": "Vocabulary",
    "grammar": "Grammar & Sentence Patterns",
    "writing": "Writing & Characters",
    "exercises": "Exercises & Homework",
    "resources": "Resources & Links",
    "other": "Additional Content",
}

CATEGORY_ORDER = ["overview", "phonetics", "grammar", "writing", "vocabulary",
                  "exercises", "resources", "other"]


def convert_pptx_to_md(filepath: str) -> str:
    """Convert a single PPTX file to a structured Markdown string."""
    prs = Presentation(filepath)
    filename = os.path.basename(filepath)

    # Extract all slides
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        sd = extract_slide(slide, i)
        if sd["texts"] or sd["tables"]:
            slides_data.append(sd)

    if not slides_data:
        return ""

    # Detect lesson info
    first_texts = slides_data[0]["texts"] if slides_data else []
    lesson_num = detect_lesson_number(filename, first_texts)
    lesson_title = detect_lesson_title(first_texts)

    # Categorize
    categories = categorize_content(slides_data)

    # Build markdown
    md_lines = []
    md_lines.append(f"# Lesson {lesson_num}: {lesson_title}")
    md_lines.append("")
    version = get_version()
    md_lines.append(f"**Source file:** `{filename}`")
    md_lines.append(f"**Total slides:** {len(prs.slides)}")
    md_lines.append(f"**Generated:** {version}")
    md_lines.append("")
    md_lines.append("---")
    md_lines.append("")

    for cat_key in CATEGORY_ORDER:
        cat_slides = categories.get(cat_key, [])
        if not cat_slides:
            continue

        md_lines.append(f"## {CATEGORY_TITLES[cat_key]}")
        md_lines.append("")

        for sd in cat_slides:
            content = render_slide_content(sd, heading_level=3)
            if content.strip():
                md_lines.append(content)

    return "\n".join(md_lines)


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    pptx_files = sorted(
        f for f in os.listdir(INPUT_DIR)
        if f.lower().endswith('.pptx') and not f.startswith('~')
    )

    if not pptx_files:
        print("No .pptx files found.")
        return

    print(f"Found {len(pptx_files)} PowerPoint file(s):\n")

    for f in pptx_files:
        filepath = os.path.join(INPUT_DIR, f)
        print(f"  Converting: {f}")

        md_content = convert_pptx_to_md(filepath)
        if not md_content:
            print(f"    -> Skipped (no content)")
            continue

        # Determine lesson number for output filename
        first_texts = []
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            first_texts.append(t)
            break

        lesson_num = detect_lesson_number(f, first_texts)
        out_name = f"lesson_{lesson_num}.md"
        out_path = os.path.join(OUTPUT_DIR, out_name)

        with open(out_path, "w", encoding="utf-8") as fh:
            fh.write(md_content)

        print(f"    -> Saved: {out_name}")

    print(f"\nDone! Markdown files saved to: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
